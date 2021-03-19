# ===================================================================================
import sys
sys.path.append("../MyLogger/")
sys.path.append("../MyDataBase/")
sys.path.append('../MyBaseSystem/')
from MyLogger import MyLogger
MyLogger = MyLogger.GetInstance()
from MyDataBase import MyDataBase
from MyTkRoot import MyTkRoot
from WidgetFactory import WidgetFactory
from MyBaseSystem import *
# ===================================================================================
import tkinter as tk
import glob
import os
# ===================================================================================
class ExplorerFrame(tk.Frame):
    @MyLogger.deco
    def __init__(self,master, cnf={},**kw):
        super().__init__(master,cnf,**kw)
        self.explorerdata = MyDataBase("explorer.xlsx")
        self.explorerdata.DBAppendColumn(['abspath', 'dirname', 'basename', 'extension', 'contents'])
        self.combobox = {}
        if 'id' in self.combobox:
            WidgetFactory.Destroy(self.combobox['id'])
        self.combobox = WidgetFactory.NewMultiCombobox(self, ['combobox'], 0,0,1,0.1, "ToBottom")
# ===================================================================================
    # @note
    # 今は隠しファイルヒットしないかも
    @MyLogger.deco
    def __Glob(self, basedir):
        # ===================================================================================
        def __xlsx2txt(path):
            ret = ""
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path)
                for ws in wb:
                    for row in ws:
                        for cell in row:
                            if cell.value == None:
                                ret += "\t"
                            else:
                                ret += str(cell.value) + "\t"
                        ret += "\n"
            except Exception as e:
                MyLogger.error(e)
            return ret
        # ===================================================================================
        def __docx2txt(path):
            ret = ""
            try:
                import docx
                doc= docx.Document(path)
                for par in doc.paragraphs:
                    ret += par.text + "\n"
            except Exception as e:
                MyLogger.error(e)
                ret = "[could not read contents]" + str(e)
            return ret
        # ===================================================================================
        def __pptx2txt(path):
            ret = ""
            try:
                import pptx
                prs = pptx.Presentation(path)
                for i, sld in enumerate(prs.slides, start=1):
                    for shp in sld.shapes:
                        if shp.has_text_frame:
                            print(shp.text)
                            ret += shp.text + "\n"
            except Exception as e:
                MyLogger.error(e)
                ret = "[could not read contents]" + str(e)
            return ret
        # ===================================================================================
        def __pdf2txt(path):
            ret = ""
            try:
                from pdfminer.high_level import extract_text
                ret = extract_text(path)
            except Exception as e:
                MyLogger.error(e)
                ret = "[could not read contents]" + str(e)
            return ret
        # ===================================================================================
        def __other2txt(path):
            ret = ""
            try:
                ret = "".join(myopen(path, 'r').readlines())
            except Exception as e:
                MyLogger.error(e)
                ret = "[could not read contents]" + str(e)
            return ret
        # ===================================================================================
        def __smallglob(basedir):
            files = []
            dirs = [basedir]
            while dirs:
                MyLogger.sakura("remaining dirs",len(dirs))
                basedir = dirs.pop(0)
                try:
                    files += [file for file in glob.glob(basedir + "/*", recursive=False) if os.path.isfile(file)]
                    dirs += [file for file in glob.glob(basedir + "/*", recursive=False) if os.path.isdir(file)]
                except Exception as e:
                    MyLogger.error("[glob failed at",basedir,"]",e)
            return files
        # path整理
        basedir = basedir.replace("\"", "")
        basedir = basedir.replace("/", os.sep)
        basedir = basedir.replace("\\", os.sep)
        MyLogger.sakura(basedir)
        if not os.path.exists(basedir) or not os.path.isdir(basedir):
            return
        # files = [file for file in glob.glob(basedir + "/**/*", recursive=True) if os.path.isfile(file)]
        files = __smallglob(basedir)
        self.explorerdata.DBRead()
        temp_dict = self.explorerdata.GetDict()
        MyLogger.SetFraction(len(files))
        for file in files:
            MyLogger.SetNumerator(files.index(file))
            MyLogger.sakura(file)
            abspath = os.path.abspath(file)
            dirname = os.path.dirname(file)
            basename = os.path.basename(file)
            extension = os.path.splitext(file)[1]
            contents = ""
            if extension == ".xlsx":
                contents = __xlsx2txt(abspath)
            elif extension == ".docx":
                contents = __docx2txt(abspath)
            elif extension == ".pptx":
                contents = __pptx2txt(abspath)
            elif extension == ".pdf":
                contents = __pdf2txt(abspath)
            else:
                contents = __other2txt(abspath)
            temp_dict = self.explorerdata.DBAppendRow([abspath, dirname, basename, extension, contents], temp_dict)
        self.explorerdata.DBImportDict(temp_dict)
# ===================================================================================
    @MyLogger.deco
    def OnKeyEvent(self, event):
        if event.keysym == 'Return':
            if (focused := WidgetFactory.HasFocus(self.combobox['id'] , self.master.focus_get())):
                path = self.combobox['instance'].comboboxes[focused]['instance'].GetText()
                self.__Glob(path)
# ===================================================================================
if __name__ == '__main__':
    root = MyTkRoot()
    explorerframe = ExplorerFrame(root)
    root.AddFrame(explorerframe, 'explorerframe', key=explorerframe.OnKeyEvent)
    root.mainloop()
# ===================================================================================